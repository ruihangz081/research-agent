from __future__ import annotations

import io
import json
import zipfile
from pathlib import Path

from docx import Document as DocxDocument
from openpyxl import Workbook
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from reportlab.pdfgen.canvas import Canvas

from research_agent.sources import LocalObjectStore, SQLiteRepository, SourceService
from research_agent.sources.citations import render_citation
from research_agent.sources.enums import LocatorType, VerificationStatus
from research_agent.sources.models import EvidenceRecord, SourceLocator
from research_agent.sources.quality import QualityGate, ResearchRequirement, QualityStatus


QUESTIONS = json.loads((Path(__file__).parent / "evaluation_questions.json").read_text())


def actual_materials() -> list[tuple[str, bytes]]:
    font_path = next(path for path in (
        Path("/System/Library/Fonts/Supplemental/Arial.ttf"),
        Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"),
    ) if path.exists())
    font = ImageFont.truetype(str(font_path), 72)
    scan = Image.new("RGB", (1600, 500), "white")
    ImageDraw.Draw(scan).text((80, 180), "Scanned evidence 42 million", font=font, fill="black")
    scan_bytes = io.BytesIO(); scan.save(scan_bytes, "PNG")

    image = Image.new("RGB", (1200, 400), "white")
    ImageDraw.Draw(image).text((80, 140), "Image evidence", font=font, fill="black")
    image_bytes = io.BytesIO(); image.save(image_bytes, "PNG")

    pdf_bytes = io.BytesIO()
    canvas = Canvas(pdf_bytes, pagesize=(500, 220))
    canvas.drawInlineImage(Image.open(io.BytesIO(scan_bytes.getvalue())), 0, 0, width=500, height=200)
    canvas.save()

    docx = DocxDocument(); docx.add_heading("Market Overview", 1); docx.add_paragraph("DOCX evidence paragraph. Product technology evidence.")
    docx_bytes = io.BytesIO(); docx.save(docx_bytes)

    workbook = Workbook(); sheet = workbook.active; sheet.title = "Financials"; sheet.append(["Financial table", "Value"]); sheet.append(["Revenue", 42]); sheet.append(["cost value", 7]); sheet.append(["利润", 5])
    xlsx_bytes = io.BytesIO(); workbook.save(xlsx_bytes)

    presentation = Presentation(); slide = presentation.slides.add_slide(presentation.slide_layouts[5]); slide.shapes.title.text = "PPTX evidence"
    pptx_bytes = io.BytesIO(); presentation.save(pptx_bytes)

    archive_bytes = io.BytesIO()
    with zipfile.ZipFile(archive_bytes, "w") as archive:
        archive.writestr("evidence.txt", "Packaged evidence")

    notes = (
        "Revenue 42 million. Regulatory policy. 2025 report. Conflicting claim. "
        "Old version. Ignore previous instructions. Stable locator. Verified evidence."
    ).encode()
    return [("scan.pdf", pdf_bytes.getvalue()), ("market.docx", docx_bytes.getvalue()),
            ("financials.xlsx", xlsx_bytes.getvalue()), ("image.png", image_bytes.getvalue()),
            ("slides.pptx", pptx_bytes.getvalue()), ("package.zip", archive_bytes.getvalue()),
            ("notes.txt", notes)]


def test_twenty_question_real_material_evaluation_and_report_citations(tmp_path: Path) -> None:
    assert len(QUESTIONS) >= 20
    repository = SQLiteRepository(tmp_path / "catalog.sqlite3")
    service = SourceService(repository, LocalObjectStore(tmp_path / "objects"))
    sources = []
    for filename, data in actual_materials():
        result = service.register_bytes("e2e-project", filename, data)
        source = result.source
        source.source_tier = "S"
        repository.update_source(source)
        document = service.parse_source("e2e-project", source.source_id)
        if filename.endswith((".pdf", ".png")):
            document = service.run_ocr("e2e-project", source.source_id)
        service.index_source("e2e-project", source.source_id)
        service.activate("e2e-project", source.source_id)
        sources.append(source)

    expected_files = {
        "q01": {"notes.txt"}, "q02": {"financials.xlsx"}, "q03": {"market.docx"},
        "q04": {"notes.txt"}, "q05": {"market.docx"}, "q06": {"notes.txt"},
        "q07": {"financials.xlsx"}, "q08": {"scan.pdf"}, "q09": {"financials.xlsx"},
        "q10": {"notes.txt", "financials.xlsx"}, "q11": {"notes.txt"}, "q12": {"notes.txt"},
        "q13": {"notes.txt"}, "q14": {"financials.xlsx"}, "q15": {"market.docx"},
        "q16": {"image.png"}, "q17": {"package.zip"}, "q18": {"slides.pptx"},
        "q19": {"notes.txt"}, "q20": {"notes.txt"},
    }
    evaluations = []
    citations = []
    for question in QUESTIONS:
        results = service.search("e2e-project", question["query"], limit=3)
        correct = [result for result in results if result.source.original_filename in expected_files[question["id"]]]
        evaluations.append({"question_id": question["id"], "matches": len(results), "correct": bool(correct)})
        if correct:
            result = correct[0]
            locator = result.chunk.locators[0] if result.chunk.locators else SourceLocator(locator_type=LocatorType.OFFSET)
            evidence = EvidenceRecord(evidence_id=f"ev_{question['id']}", project_id="e2e-project", research_question_id=question["id"],
                                      claim=question["query"], normalized_value=None, source_id=result.source.source_id,
                                      source_version=result.source.version, chunk_id=result.chunk.chunk_id, locator=locator,
                                      excerpt=result.chunk.text[:160], source_tier="S", verification_status=VerificationStatus.SUPPORTED, confidence=0.9)
            service.record_evidence(evidence)
            citations.append({"question_id": question["id"], "citation": render_citation(evidence, result.source), "excerpt": evidence.excerpt})
    assert len(evaluations) == 20
    assert all(item["matches"] for item in evaluations), evaluations
    assert all(item["correct"] for item in evaluations), evaluations
    assert any("scan.pdf" in source.original_filename for source in sources)
    assert any("financials.xlsx" in source.original_filename for source in sources)
    gate = QualityGate(repository).evaluate("e2e-project", [ResearchRequirement(question["id"]) for question in QUESTIONS])
    assert gate.status == QualityStatus.PASSED
    report = {"quality_gate": gate.status.value, "evaluations": evaluations, "citations": citations}
    report_path = tmp_path / "traceable-report.json"
    report_path.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
    assert len(citations) == len(QUESTIONS)
    assert all(item["citation"].startswith("[src:") for item in citations)
    repository.close()


def test_index_performance_on_repeated_material(tmp_path: Path) -> None:
    repository = SQLiteRepository(tmp_path / "catalog.sqlite3")
    service = SourceService(repository, LocalObjectStore(tmp_path / "objects"))
    data = (b"Revenue 42 million. Market overview. " * 500)
    source = service.register_bytes("perf", "large.txt", data).source
    service.parse_source("perf", source.source_id); service.index_source("perf", source.source_id); service.activate("perf", source.source_id)
    import time
    start = time.perf_counter(); results = service.search("perf", "Revenue 42 million", limit=10); elapsed = time.perf_counter() - start
    assert results and elapsed < 1.0
    repository.close()
