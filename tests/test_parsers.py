from __future__ import annotations

import io
import json
import zipfile
from pathlib import Path

from docx import Document as DocxDocument
from openpyxl import Workbook
from PIL import Image
from pptx import Presentation
from reportlab.pdfgen.canvas import Canvas

from research_agent.sources import SourceService
from research_agent.sources.parsers import parse_bytes
from research_agent.sources.enums import LocatorType


EXPECTED = json.loads((Path(__file__).parent / "golden/expected.json").read_text())


def fixture_bytes(kind: str) -> tuple[str, bytes]:
    if kind == "txt": return "evidence.txt", b"Revenue reached 42 million"
    if kind == "md": return "evidence.md", b"# Market Overview\nA stable paragraph"
    if kind == "html": return "evidence.html", b"<html><title>Golden HTML</title><p>Verified web evidence</p></html>"
    if kind == "mhtml":
        return "evidence.mhtml", b"MIME-Version: 1.0\nContent-Type: multipart/related; boundary=abc\n\n--abc\nContent-Type: text/html\n\n<html><p>Archived web evidence</p></html>\n--abc--"
    if kind in {"csv", "tsv"}:
        delimiter = "," if kind == "csv" else "\t"
        row = f"Revenue{delimiter}42" if kind == "csv" else f"利润{delimiter}42"
        return f"evidence.{kind}", f"Metric{delimiter}Value\n{row}".encode()
    if kind == "rtf": return "evidence.rtf", b"{\\rtf1\\ansi RTF evidence}"
    if kind == "image":
        output = io.BytesIO()
        Image.new("RGB", (120, 60), "white").save(output, format="PNG")
        return "evidence.png", output.getvalue()
    if kind == "pdf":
        output = io.BytesIO()
        canvas = Canvas(output)
        canvas.drawString(72, 720, "Golden PDF evidence")
        canvas.showPage()
        canvas.save()
        return "evidence.pdf", output.getvalue()
    if kind == "docx":
        document = DocxDocument()
        document.add_paragraph("DOCX evidence paragraph")
        output = io.BytesIO()
        document.save(output)
        return "evidence.docx", output.getvalue()
    if kind == "xlsx":
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "Financials"
        sheet.append(["Metric", "Value"])
        sheet.append(["Revenue", 42])
        output = io.BytesIO()
        workbook.save(output)
        return "evidence.xlsx", output.getvalue()
    if kind == "pptx":
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "PPTX evidence slide"
        output = io.BytesIO()
        presentation.save(output)
        return "evidence.pptx", output.getvalue()
    if kind == "zip":
        output = io.BytesIO()
        with zipfile.ZipFile(output, "w") as archive:
            archive.writestr("inside.txt", "Packaged evidence")
        return "evidence.zip", output.getvalue()
    raise AssertionError(kind)


def test_all_golden_formats_have_expected_content_and_locator() -> None:
    for kind in EXPECTED:
        filename, data = fixture_bytes(kind)
        document = parse_bytes(f"golden_{kind}", data, filename).document
        expected = EXPECTED[kind]
        text = "\n".join(block.text for block in document.blocks)
        if "text" in expected:
            assert expected["text"] in text, kind
        if "page" in expected:
            assert document.pages[0].page_number == expected["page"]
            assert document.blocks[0].locator.locator_type == LocatorType.PAGE
        if "sheet" in expected:
            assert document.tables[0].sheet_name == expected["sheet"]
            assert next(cell for cell in document.tables[0].cells if cell.row == 2 and cell.column == 2).value == "42"
            assert document.tables[0].locator.locator_type == LocatorType.SHEET
        if "slide" in expected:
            assert document.pages[0].page_number == expected["slide"]
            assert document.blocks[0].locator.locator_type == LocatorType.SLIDE
        if kind == "zip":
            assert document.blocks[0].locator.locator_type == LocatorType.ZIP_MEMBER
            assert document.blocks[0].locator.zip_member == expected["member"]
        if kind == "image":
            assert document.images[0].bbox == (0, 0, expected["width"], expected["height"])


def test_parser_output_is_stored_as_derived_layer(tmp_path: Path) -> None:
    from research_agent.sources import LocalObjectStore, SQLiteRepository
    repository = SQLiteRepository(tmp_path / "catalog.sqlite3")
    service = SourceService(repository, LocalObjectStore(tmp_path / "objects"))
    filename, data = fixture_bytes("pdf")
    result = service.register_bytes("project", filename, data)
    document = service.parse_source("project", result.source.source_id)
    assert document.source_id == result.source.source_id
    assert service.raw_bytes("project", result.source.source_id) == data
    assert service.get_source("project", result.source.source_id).status.value == "ready"
    assert repository.get_document(result.source.source_id, "project") is not None
    repository.close()


def test_text_duplicate_lines_keep_distinct_offsets() -> None:
    document = parse_bytes("text", b"alpha\nbeta\nalpha", "repeat.txt").document
    assert [(block.locator.char_start, block.locator.char_end) for block in document.blocks] == [(0, 5), (6, 10), (11, 16)]


def test_csv_range_supports_more_than_twenty_six_columns() -> None:
    document = parse_bytes("wide", (",".join(str(value) for value in range(30))).encode(), "wide.csv").document
    assert document.tables[0].locator.cell_range == "A1:AD1"
