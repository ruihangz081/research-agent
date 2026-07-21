"""Format-aware parsers that normalize material into SourceDocument."""
from __future__ import annotations

import csv
import io
import mimetypes
import re
import shutil
import subprocess
import tempfile
import zipfile
from email import policy
from email.parser import BytesParser
from pathlib import Path
from typing import Callable

from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from openpyxl import load_workbook
from openpyxl.utils import get_column_letter
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pypdf import PdfReader
from striprtf.striprtf import rtf_to_text

from ..enums import BlockType, LocatorType
from ..models import (
    ContentBlock, DocumentMetadata, ExtractionQuality, ExtractionWarning, ImageBlock,
    PageInfo, SourceDocument, SourceLocator, TableBlock, TableCell,
)
from ..security import inspect_zip, sanitize_untrusted_text


class ParseError(ValueError):
    pass


class ParseResult:
    def __init__(self, document: SourceDocument):
        self.document = document


def _decode(data: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "gb18030", "utf-16"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _block(source_id: str, index: int, text: str, block_type: BlockType, locator: SourceLocator | None = None, **kwargs) -> ContentBlock:
    text, injection_warnings = sanitize_untrusted_text(text)
    attrs = dict(kwargs.pop("attributes", {}))
    if injection_warnings:
        attrs["untrusted_instruction_warning"] = injection_warnings
    return ContentBlock(block_id=f"blk_{source_id}_{index}", source_id=source_id, block_type=block_type,
                        text=text, order=index, locator=locator, attributes=attrs, **kwargs)


def _quality(pages: list[PageInfo], blocks: list[ContentBlock], tables: list[TableBlock], warnings: list[ExtractionWarning]) -> ExtractionQuality:
    text_blocks = sum(bool(block.text.strip()) for block in blocks)
    text_coverage = (sum(bool(page.text.strip()) for page in pages) / len(pages)) if pages else float(bool(text_blocks))
    located = sum(bool(block.locator) for block in blocks)
    layout_coverage = located / len(blocks) if blocks else 0
    usable_tables = sum(bool(table.cells and table.locator) for table in tables)
    table_coverage = usable_tables / len(tables) if tables else 0
    penalty = min(0.5, sum(0.25 if item.severity in {"high", "error"} else 0.05 for item in warnings))
    score = max(0.0, min(1.0, 0.55 * text_coverage + 0.3 * layout_coverage + 0.15 * table_coverage - penalty))
    return ExtractionQuality(score=score, text_coverage=text_coverage,
                             layout_coverage=layout_coverage, table_coverage=table_coverage,
                             warnings=len(warnings))


def _document(source_id: str, blocks: list[ContentBlock], *, pages=None, tables=None, images=None, warnings=None, metadata=None) -> SourceDocument:
    pages, tables, images, warnings = pages or [], tables or [], images or [], warnings or []
    return SourceDocument(source_id=source_id, document_id=f"doc_{source_id}", metadata=metadata or DocumentMetadata(),
                          pages=pages, blocks=blocks, tables=tables, images=images,
                          warnings=warnings, quality=_quality(pages, blocks, tables, warnings))


def parse_text(source_id: str, data: bytes, filename: str) -> SourceDocument:
    text = _decode(data)
    blocks: list[ContentBlock] = []
    headings: list[str] = []
    cursor = 0
    for index, raw_line in enumerate(text.splitlines(keepends=True)):
        line = raw_line.rstrip("\r\n")
        stripped = line.strip()
        if not stripped:
            cursor += len(raw_line)
            continue
        if filename.lower().endswith((".md", ".markdown")) and stripped.startswith("#"):
            level = len(stripped) - len(stripped.lstrip("#"))
            headings = headings[: level - 1] + [stripped[level:].strip()]
            block_type = BlockType.TITLE if level == 1 and index == 0 else BlockType.HEADING
        else:
            block_type = BlockType.PARAGRAPH
        char_start = cursor + len(line) - len(line.lstrip())
        blocks.append(_block(source_id, len(blocks), stripped, block_type,
                             SourceLocator(locator_type=LocatorType.OFFSET, char_start=char_start, char_end=char_start + len(stripped)),
                             heading_path=headings.copy()))
        cursor += len(raw_line)
    return _document(source_id, blocks)


def parse_html(source_id: str, data: bytes, filename: str) -> SourceDocument:
    if filename.lower().endswith(".mhtml"):
        message = BytesParser(policy=policy.default).parsebytes(data)
        parts = [part for part in message.walk() if part.get_content_type() == "text/html"]
        html = parts[0].get_content() if parts else _decode(data)
    else:
        html = _decode(data)
    soup = BeautifulSoup(html, "html.parser")
    blocks: list[ContentBlock] = []
    headings: list[str] = []
    for element in soup.find_all(["title", "h1", "h2", "h3", "h4", "p", "li", "blockquote", "pre"]):
        text = " ".join(element.get_text(" ", strip=True).split())
        if not text:
            continue
        tag = element.name
        kind = BlockType.TITLE if tag == "title" else BlockType.HEADING if tag.startswith("h") else BlockType.LIST if tag == "li" else BlockType.QUOTE if tag == "blockquote" else BlockType.PARAGRAPH
        if tag.startswith("h"):
            level = int(tag[1])
            headings = headings[: level - 1] + [text]
        locator = SourceLocator(
            locator_type=LocatorType.PARAGRAPH,
            paragraph_index=len(blocks),
            heading_path=headings.copy(),
        )
        blocks.append(_block(source_id, len(blocks), text, kind, locator, heading_path=headings.copy()))
    tables: list[TableBlock] = []
    for table_index, raw_table in enumerate(soup.find_all("table"), 1):
        rows = [[cell.get_text(" ", strip=True) for cell in row.find_all(["th", "td"])] for row in raw_table.find_all("tr")]
        cells = [TableCell(row=r + 1, column=c + 1, value=value) for r, row in enumerate(rows) for c, value in enumerate(row)]
        table_id = f"tbl_{source_id}_{table_index}"
        locator = SourceLocator(locator_type=LocatorType.RANGE, table_id=table_id)
        tables.append(TableBlock(table_id=table_id, source_id=source_id, rows=len(rows),
                                 columns=max((len(row) for row in rows), default=0), cells=cells, locator=locator))
        blocks.append(_block(source_id, len(blocks), "\n".join(" | ".join(row) for row in rows), BlockType.TABLE, locator))
    author = soup.find("meta", attrs={"name": re.compile("author", re.I)})
    publisher = soup.find("meta", attrs={"name": re.compile("publisher|site_name", re.I)})
    return _document(source_id, blocks, tables=tables, metadata=DocumentMetadata(
        title=soup.title.get_text(strip=True) if soup.title else None,
        authors=[author.get("content")] if author and author.get("content") else [],
        publisher=publisher.get("content") if publisher else None,
    ))


def parse_csv(source_id: str, data: bytes, filename: str) -> SourceDocument:
    delimiter = "\t" if filename.lower().endswith(".tsv") else ","
    rows = list(csv.reader(io.StringIO(_decode(data)), delimiter=delimiter))
    cells = [TableCell(row=r + 1, column=c + 1, value=value) for r, row in enumerate(rows) for c, value in enumerate(row)]
    columns = max((len(row) for row in rows), default=1)
    last_row = max(len(rows), 1)
    table = TableBlock(table_id=f"tbl_{source_id}_1", source_id=source_id, rows=len(rows), columns=columns, cells=cells,
                       range=f"A1:{get_column_letter(columns)}{last_row}",
                       locator=SourceLocator(locator_type=LocatorType.RANGE, table_id=f"tbl_{source_id}_1", cell_range=f"A1:{get_column_letter(columns)}{last_row}"))
    text = "\n".join(" | ".join(row) for row in rows)
    block = _block(source_id, 0, text, BlockType.TABLE, table.locator)
    return _document(source_id, [block], tables=[table])


def parse_pdf(source_id: str, data: bytes, filename: str) -> SourceDocument:
    try:
        reader = PdfReader(io.BytesIO(data))
    except Exception as exc:
        raise ParseError(f"cannot read PDF: {exc}") from exc
    if reader.is_encrypted:
        try:
            if not reader.decrypt(""):
                return _document(source_id, [], warnings=[ExtractionWarning(code="encrypted_pdf", message="PDF password required", severity="high", method="pypdf")])
        except Exception:
            return _document(source_id, [], warnings=[ExtractionWarning(code="encrypted_pdf", message="PDF password required", severity="high", method="pypdf")])
    import fitz

    pages, blocks, tables, images, warnings = [], [], [], [], []
    document = fitz.open(stream=data, filetype="pdf")
    for page_number, page in enumerate(document, 1):
        page_blocks = page.get_text("blocks", sort=True)
        text_parts: list[str] = []
        for raw_block in page_blocks:
            x0, y0, x1, y1, raw_text = raw_block[:5]
            block_text = " ".join(str(raw_text).split())
            if not block_text:
                continue
            text_parts.append(block_text)
            blocks.append(_block(
                source_id, len(blocks), block_text, BlockType.PARAGRAPH,
                SourceLocator(locator_type=LocatorType.PAGE, page_number=page_number, bbox=(x0, y0, x1, y1)),
                page_number=page_number,
            ))
        text = "\n".join(text_parts)
        pages.append(PageInfo(page_number=page_number, width=page.rect.width, height=page.rect.height,
                              rotation=page.rotation, text=text, is_scanned=not bool(text)))
        if not text:
            warnings.append(ExtractionWarning(code="scanned_page", message="page has no native text; OCR required", page_number=page_number, method="pymupdf"))
        try:
            found_tables = page.find_tables().tables
        except Exception:
            found_tables = []
        for table_index, raw_table in enumerate(found_tables, 1):
            extracted = raw_table.extract()
            cells = [TableCell(row=r + 1, column=c + 1, value="" if value is None else str(value))
                     for r, row in enumerate(extracted) for c, value in enumerate(row)]
            table_id = f"tbl_{source_id}_{page_number}_{table_index}"
            bbox = tuple(float(value) for value in raw_table.bbox)
            locator = SourceLocator(locator_type=LocatorType.PAGE, page_number=page_number, table_id=table_id, bbox=bbox)
            tables.append(TableBlock(table_id=table_id, source_id=source_id, page_number=page_number,
                                     rows=len(extracted), columns=max((len(row) for row in extracted), default=0),
                                     cells=cells, locator=locator))
        for image_index, raw_image in enumerate(page.get_images(full=True), 1):
            xref = raw_image[0]
            rects = page.get_image_rects(xref)
            for rect_index, rect in enumerate(rects, 1):
                images.append(ImageBlock(image_id=f"img_{source_id}_{page_number}_{image_index}_{rect_index}",
                                         source_id=source_id, page_number=page_number,
                                         bbox=(rect.x0, rect.y0, rect.x1, rect.y1)))
    metadata = DocumentMetadata(title=(reader.metadata.title if reader.metadata else None), authors=[reader.metadata.author] if reader.metadata and reader.metadata.author else [])
    return _document(source_id, blocks, pages=pages, tables=tables, images=images, warnings=warnings, metadata=metadata)


def parse_docx(source_id: str, data: bytes, filename: str) -> SourceDocument:
    document = DocxDocument(io.BytesIO(data))
    blocks: list[ContentBlock] = []
    headings: list[str] = []
    for paragraph_index, paragraph in enumerate(document.paragraphs):
        text = paragraph.text.strip()
        if text:
            style = paragraph.style.name.lower() if paragraph.style else ""
            kind = BlockType.TITLE if "title" in style else BlockType.HEADING if "heading" in style else BlockType.PARAGRAPH
            match = re.search(r"heading\s*(\d+)", style)
            if match:
                level = int(match.group(1))
                headings = headings[:level - 1] + [text]
            blocks.append(_block(source_id, len(blocks), text, kind,
                                 SourceLocator(locator_type=LocatorType.PARAGRAPH, paragraph_index=paragraph_index, heading_path=headings.copy()),
                                 heading_path=headings.copy()))
    tables: list[TableBlock] = []
    for table_index, raw_table in enumerate(document.tables, 1):
        cells = [TableCell(row=r + 1, column=c + 1, value=cell.text.strip()) for r, row in enumerate(raw_table.rows) for c, cell in enumerate(row.cells)]
        table = TableBlock(table_id=f"tbl_{source_id}_{table_index}", source_id=source_id, rows=len(raw_table.rows), columns=len(raw_table.columns), cells=cells,
                           locator=SourceLocator(locator_type=LocatorType.RANGE, table_id=f"tbl_{source_id}_{table_index}"))
        tables.append(table)
        blocks.append(_block(source_id, len(blocks), "\n".join(" | ".join(cell.value for cell in row) for row in [[c for c in cells if c.row == r] for r in range(1, table.rows + 1)]), BlockType.TABLE, table.locator))
    images = [ImageBlock(image_id=f"img_{source_id}_{index}", source_id=source_id,
                         bbox=(0, 0, float(shape.width), float(shape.height)))
              for index, shape in enumerate(document.inline_shapes, 1)]
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as package:
            for member, kind in (("word/footnotes.xml", BlockType.FOOTNOTE), ("word/endnotes.xml", BlockType.FOOTNOTE), ("word/comments.xml", BlockType.QUOTE)):
                if member not in package.namelist():
                    continue
                xml = BeautifulSoup(package.read(member), "xml")
                for item in xml.find_all(["w:footnote", "w:endnote", "w:comment", "footnote", "endnote", "comment"]):
                    text = " ".join(value.get_text(" ", strip=True) for value in item.find_all(["w:t", "t"]))
                    if text:
                        blocks.append(_block(source_id, len(blocks), text, kind,
                                             SourceLocator(locator_type=LocatorType.PARAGRAPH, paragraph_index=len(blocks))))
    except zipfile.BadZipFile:
        pass
    return _document(source_id, blocks, tables=tables, images=images)


def parse_xlsx(source_id: str, data: bytes, filename: str) -> SourceDocument:
    workbook = load_workbook(io.BytesIO(data), read_only=False, data_only=False)
    values_workbook = load_workbook(io.BytesIO(data), read_only=False, data_only=True)
    blocks, tables = [], []
    for sheet in workbook.worksheets:
        rows = list(sheet.iter_rows(values_only=False))
        value_sheet = values_workbook[sheet.title]
        cells = []
        for r, row in enumerate(rows, 1):
            for c, cell in enumerate(row, 1):
                formula = str(cell.value) if isinstance(cell.value, str) and cell.value.startswith("=") else None
                displayed = value_sheet.cell(r, c).value if formula else cell.value
                cells.append(TableCell(row=r, column=c, value="" if displayed is None else str(displayed), formula=formula))
        table_id = f"tbl_{source_id}_{sheet.title}"
        columns = max((len(row) for row in rows), default=1)
        cell_range = f"A1:{get_column_letter(columns)}{max(len(rows), 1)}"
        table = TableBlock(table_id=table_id, source_id=source_id, sheet_name=sheet.title, rows=len(rows), columns=max((len(row) for row in rows), default=0), cells=cells,
                           range=cell_range,
                           locator=SourceLocator(locator_type=LocatorType.SHEET, sheet_name=sheet.title, table_id=table_id, cell_range=cell_range))
        tables.append(table)
        text = sheet.title + "\n" + "\n".join(" | ".join(c.value for c in cells if c.row == r) for r in range(1, table.rows + 1))
        blocks.append(_block(source_id, len(blocks), text, BlockType.TABLE, table.locator, attributes={"sheet_name": sheet.title}))
        blocks[-1].attributes.update({
            "sheet_state": sheet.sheet_state,
            "merged_ranges": [str(value) for value in sheet.merged_cells.ranges],
            "hidden_rows": [index for index, value in sheet.row_dimensions.items() if value.hidden],
            "hidden_columns": [index for index, value in sheet.column_dimensions.items() if value.hidden],
        })
    return _document(source_id, blocks, tables=tables)


def parse_pptx(source_id: str, data: bytes, filename: str) -> SourceDocument:
    presentation = Presentation(io.BytesIO(data))
    blocks, pages, tables, images = [], [], [], []
    for slide_number, slide in enumerate(presentation.slides, 1):
        texts = []
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                cells = [TableCell(row=r + 1, column=c + 1, value=value) for r, row in enumerate(rows) for c, value in enumerate(row)]
                table_id = f"tbl_{source_id}_{slide_number}_{len(tables) + 1}"
                locator = SourceLocator(locator_type=LocatorType.SLIDE, slide_number=slide_number, table_id=table_id,
                                        bbox=(float(shape.left), float(shape.top), float(shape.left + shape.width), float(shape.top + shape.height)))
                tables.append(TableBlock(table_id=table_id, source_id=source_id, page_number=slide_number,
                                         rows=len(rows), columns=max((len(row) for row in rows), default=0), cells=cells, locator=locator))
                texts.append("\n".join(" | ".join(row) for row in rows))
                continue
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                images.append(ImageBlock(image_id=f"img_{source_id}_{slide_number}_{len(images) + 1}", source_id=source_id,
                                         page_number=slide_number,
                                         bbox=(float(shape.left), float(shape.top), float(shape.left + shape.width), float(shape.top + shape.height))))
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
            if getattr(shape, "has_chart", False):
                chart_title = shape.chart.chart_title.text_frame.text if shape.chart.has_title else "chart"
                series = [getattr(item, "name", "series") for item in shape.chart.series]
                texts.append(f"{chart_title}: {', '.join(str(value) for value in series)}")
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                texts.append(f"Speaker notes: {notes}")
        text = "\n".join(texts)
        pages.append(PageInfo(page_number=slide_number, text=text))
        if text:
            blocks.append(_block(source_id, len(blocks), text, BlockType.PARAGRAPH, SourceLocator(locator_type=LocatorType.SLIDE, slide_number=slide_number)))
    return _document(source_id, blocks, pages=pages, tables=tables, images=images)


def parse_image(source_id: str, data: bytes, filename: str) -> SourceDocument:
    try:
        image = Image.open(io.BytesIO(data))
        image.verify()
        width, height = image.size
    except Exception as exc:
        raise ParseError(f"invalid image: {exc}") from exc
    image_block = ImageBlock(image_id=f"img_{source_id}_1", source_id=source_id, page_number=1, bbox=(0, 0, float(width), float(height)))
    warning = ExtractionWarning(code="ocr_pending", message="image text requires OCR worker", severity="warning", method="ocr")
    return _document(source_id, [], images=[image_block], warnings=[warning])


def parse_rtf(source_id: str, data: bytes, filename: str) -> SourceDocument:
    return parse_text(source_id, rtf_to_text(_decode(data)).encode(), filename)


def parse_archive(source_id: str, data: bytes, filename: str, *, depth: int = 0, max_depth: int = 3) -> SourceDocument:
    from .registry import parse_bytes
    if depth > max_depth:
        return _document(source_id, [], warnings=[ExtractionWarning(code="archive_depth_exceeded", message=f"nested archive depth exceeds {max_depth}", severity="high", method="zip")])
    inspect_zip(data)
    blocks, tables, pages, images, warnings = [], [], [], [], []
    with zipfile.ZipFile(io.BytesIO(data)) as archive:
        for member in archive.infolist():
            if member.is_dir():
                continue
            suffix = Path(member.filename).suffix.lower()
            if suffix == ".zip":
                nested = parse_archive(f"{source_id}_{len(blocks)}", archive.read(member), member.filename, depth=depth + 1, max_depth=max_depth)
                for block in nested.blocks:
                    inner = block.locator.zip_member if block.locator else None
                    block.source_id = source_id
                    block.block_id = f"blk_{source_id}_{len(blocks)}"
                    block.locator = SourceLocator(locator_type=LocatorType.ZIP_MEMBER, zip_member=f"{member.filename}/{inner}" if inner else member.filename)
                    blocks.append(block)
                tables.extend(nested.tables)
                pages.extend(nested.pages)
                images.extend(nested.images)
                warnings.extend(nested.warnings)
                continue
            try:
                nested = parse_bytes(f"{source_id}_{len(blocks)}", archive.read(member), member.filename).document
                for block in nested.blocks:
                    block.locator = SourceLocator(locator_type=LocatorType.ZIP_MEMBER, zip_member=member.filename, page_number=block.locator.page_number if block.locator else None)
                    blocks.append(block)
                tables.extend(nested.tables)
                pages.extend(nested.pages)
                images.extend(nested.images)
                warnings.extend(nested.warnings)
            except (ParseError, ValueError) as exc:
                warnings.append(ExtractionWarning(code="archive_member_failed", message=f"{member.filename}: {exc}", severity="warning", method="zip"))
    return _document(source_id, blocks, pages=pages, tables=tables, images=images, warnings=warnings)


def _legacy_convert(data: bytes, filename: str) -> tuple[bytes, str]:
    soffice = shutil.which("libreoffice") or shutil.which("soffice")
    if not soffice:
        raise ParseError("legacy Office conversion requires libreoffice")
    with tempfile.TemporaryDirectory(prefix="research-source-") as temp:
        input_path = Path(temp) / filename
        input_path.write_bytes(data)
        target = {".doc": "docx", ".xls": "xlsx", ".ppt": "pptx"}[input_path.suffix.lower()]
        subprocess.run([soffice, "--headless", "--convert-to", target, "--outdir", temp, str(input_path)], check=True, capture_output=True, timeout=60)
        converted = Path(temp) / f"{input_path.stem}.{target}"
        if not converted.exists():
            raise ParseError("libreoffice produced no converted document")
        return converted.read_bytes(), converted.name


def parse_bytes(source_id: str, data: bytes, filename: str) -> ParseResult:
    suffix = Path(filename).suffix.lower()
    if suffix in {".txt", ".md", ".markdown"}:
        document = parse_text(source_id, data, filename)
    elif suffix in {".html", ".htm", ".mhtml"}:
        document = parse_html(source_id, data, filename)
    elif suffix in {".csv", ".tsv"}:
        document = parse_csv(source_id, data, filename)
    elif suffix == ".pdf":
        document = parse_pdf(source_id, data, filename)
    elif suffix == ".docx":
        document = parse_docx(source_id, data, filename)
    elif suffix == ".xlsx":
        document = parse_xlsx(source_id, data, filename)
    elif suffix == ".pptx":
        document = parse_pptx(source_id, data, filename)
    elif suffix in {".jpg", ".jpeg", ".png", ".tif", ".tiff", ".bmp"}:
        document = parse_image(source_id, data, filename)
    elif suffix == ".rtf":
        document = parse_rtf(source_id, data, filename)
    elif suffix == ".zip":
        document = parse_archive(source_id, data, filename)
    elif suffix in {".doc", ".xls", ".ppt"}:
        converted, converted_name = _legacy_convert(data, filename)
        document = parse_bytes(source_id, converted, converted_name).document
        document.warnings.append(ExtractionWarning(code="legacy_converted", message="legacy Office file converted with LibreOffice", method="libreoffice"))
    else:
        raise ParseError(f"unsupported format: {suffix}")
    return ParseResult(document)
